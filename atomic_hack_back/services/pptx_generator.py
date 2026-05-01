"""Direct PPTX generation service using python-pptx"""
import io
import os
from http.cookiejar import debug
from logging import DEBUG
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import json
from openai import AsyncOpenAI

# Import image service for Kandinsky integration
try:
    from .image_service import generate_image_with_retry, check_kandinsky_configuration
    KANDINSKY_AVAILABLE = True
except ImportError:
    KANDINSKY_AVAILABLE = False
    print("⚠️ Image service not available, image blocks will be skipped")


class AllLayouts:
    """Обёртка для доступа ко всем layouts из всех slide masters"""
    def __init__(self, presentation):
        self.layouts = []
        for slide_master in presentation.slide_masters:
            for layout in slide_master.slide_layouts:
                self.layouts.append(layout)
    
    def __len__(self):
        return len(self.layouts)
    
    def __getitem__(self, index):
        return self.layouts[index]
    
    def __iter__(self):
        return iter(self.layouts)


def _has_usable_placeholders(layout) -> bool:
    """
    Проверяет есть ли в layout хотя бы один нормальный (не маленький) placeholder
    который мы сможем заполнить контентом.
    
    Критерии нормального placeholder:
    - Высота >= 1.0" (не подзаголовок)
    - ИЛИ: не в верхней зоне (<2.5") И высота >= 0.5"
    
    Это позволяет отфильтровать layouts где все placeholders слишком маленькие
    и не будут заполнены (например только подзаголовки/footer).
    """
    for shape in layout.placeholders:
        try:
            # Минимальный порог высоты для нормального placeholder
            min_height = Inches(1.0)
            
            # Проверяем высоту
            if hasattr(shape, 'height') and shape.height >= min_height:
                return True
            
            # Также считаем нормальным если placeholder не в верхней зоне
            # (даже если маленький - это может быть контентный placeholder)
            if hasattr(shape, 'top') and hasattr(shape, 'height'):
                is_not_in_subtitle_zone = shape.top >= Inches(2.5)
                if is_not_in_subtitle_zone and shape.height >= Inches(0.5):
                    return True
        except:
            continue
    
    return False


async def select_layout_with_ai(
    all_layouts,  # Может быть AllLayouts или prs.slide_layouts
    slide_data: Dict[str, Any],
    position: int,
    excluded_layouts: Optional[List[int]] = None
) -> int:
    """
    Использует LLM для выбора наиболее подходящего layout на основе:
    - Названий доступных layouts
    - Содержимого слайда (title, description, content_json)
    - Позиции слайда (первый = титульный)
    
    Для контентных слайдов (position > 0) отфильтровывает layouts с только
    маленькими placeholders, чтобы AI не выбирал layouts которые невозможно заполнить.
    
    Args:
        excluded_layouts: Список индексов layouts которые нужно исключить из выбора
    
    Returns:
        Индекс выбранного layout
    """
    if excluded_layouts is None:
        excluded_layouts = []
    
    # Получаем информацию о доступных layouts
    # Фильтруем layouts с только маленькими placeholders и excluded
    layouts_info = []
    for idx, layout in enumerate(all_layouts):
        # Пропускаем excluded layouts
        if idx in excluded_layouts:
            continue
        
        # Для титульного слайда (position=0) берем все layouts
        # Для остальных слайдов проверяем наличие нормальных placeholders
        if position == 0 or _has_usable_placeholders(layout):
            layouts_info.append({
                "index": idx,
                "name": layout.name,
                "placeholders": len([s for s in layout.placeholders])
            })
    
    # Логируем фильтрацию
    total_layouts = len(all_layouts)
    filtered_count = len(layouts_info)
    if position != 0 and filtered_count < total_layouts:
        print(f"🔍 Отфильтровано {total_layouts - filtered_count} layouts с маленькими placeholders ({filtered_count}/{total_layouts} осталось)")
    
    # Если отфильтровали все layouts (не должно случиться) - берем первый
    if not layouts_info:
        print("⚠️ Все layouts отфильтрованы, используем первый")
        layouts_info = [{
            "index": 0,
            "name": all_layouts[0].name,
            "placeholders": len([s for s in all_layouts[0].placeholders])
        }]
    
    # Если нет OpenAI ключа, используем простую логику
    openai_key = os.getenv("OPENAI_API_KEY")
    if not openai_key:
        # AI выбор недоступен - используем правила на основе названий layouts
        if position == 0:
            # Ищем layout с "титул" или "title" в названии
            for info in layouts_info:
                if any(word in info["name"].lower() for word in ["title", "титул", "заголов", "обложк"]):
                    return info["index"]
            # Последний fallback
            return 6 if len(all_layouts) > 6 else 0
        else:
            # Ищем layout с "контент" или "content" в названии
            for info in layouts_info:
                if any(word in info["name"].lower() for word in ["content", "контент", "текст", "слайд"]):
                    return info["index"]
            # Последний fallback
            return min(1, len(all_layouts) - 1)
    
    # Формируем запрос для LLM
    slide_title = slide_data.get('title', '')
    slide_desc = slide_data.get('description', '')
    
    # Парсим content_json для понимания типа контента
    content_json_str = slide_data.get('content_json', '{}')
    try:
        content_data = json.loads(content_json_str) if isinstance(content_json_str, str) else content_json_str
    except:
        content_data = {}
    
    blocks = content_data.get('blocks', [])
    content_types = [block.get('type', 'text') for block in blocks]
    
    # Определяем основной тип контента
    has_chart = 'chart' in content_types
    has_table = 'table' in content_types
    has_image = 'image' in content_types
    has_text = 'text' in content_types or 'list' in content_types
    
    # Формируем описание слайда
    slide_info = f"""Позиция: {position}
Описание: {slide_desc}
Типы контента: {', '.join(content_types) if content_types else 'нет контента'}"""
    
    # Формируем список layouts
    layouts_list = "\n".join([
        f"{info['index']}. {info['name']} (placeholders: {info['placeholders']})"
        for info in layouts_info
    ])
    
    # ЖЕСТКИЕ ПРАВИЛА выбора layout в зависимости от типа контента
    if position == 0:
        rules = """ЖЕСТКИЕ ПРАВИЛА ВЫБОРА:
- Это ТИТУЛЬНЫЙ слайд
- ОБЯЗАТЕЛЬНО выбирай ТОЛЬКО layout: "Титульный_белая тема" или "Титульный_белая_тема"
- Если такого нет, выбирай первый layout с "Титул" или "Title" в названии"""
    
    elif has_chart:
        rules = """ЖЕСТКИЕ ПРАВИЛА ВЫБОРА:
- На слайде есть ДИАГРАММА
- ОБЯЗАТЕЛЬНО выбирай layout: "Диаграмма с текстом" или "Диаграма с текстом"
- Если такого нет, выбирай "Title and Content" или "Слайд с заголовком"
- НЕ выбирай титульные, заключительные, картографические layouts"""
    
    elif has_image:
        rules = """ЖЕСТКИЕ ПРАВИЛА ВЫБОРА:
- На слайде есть ИЗОБРАЖЕНИЕ
- ОБЯЗАТЕЛЬНО выбирай layout: "Изображение с текстом_1" или "1 изображение"
- Если такого нет, выбирай "Title and Content" или "Слайд с заголовком"
- НЕ выбирай титульные, заключительные, картографические layouts"""
    
    elif has_table:
        rules = """ЖЕСТКИЕ ПРАВИЛА ВЫБОРА:
- На слайде есть ТАБЛИЦА
- ОБЯЗАТЕЛЬНО выбирай layout: "Таблица с текстом"
- Если такого нет, выбирай "Title and Content" или "Слайд с заголовком"
- НЕ выбирай титульные, заключительные, картографические layouts"""
    
    else:
        rules = """ЖЕСТКИЕ ПРАВИЛА ВЫБОРА:
- На слайде только ТЕКСТ/СПИСОК
- ОБЯЗАТЕЛЬНО выбирай layout: "Слайд с логотипом" или нечто подобное
- НЕ выбирай титульные, заключительные, картографические, инфографические layouts
- НЕ выбирай layouts с названиями: "Карта", "Диаграмм", "Таблиц", "Изображ", "Титул", "Заключ", "Контакт", "Выводы"
- Если "Слайд с логотипом" нет - выбирай САМЫЙ ПРОСТОЙ layout с 1-2 placeholders"""
    
    prompt = f"""Выбери наиболее подходящий layout для слайда презентации.

ИНФОРМАЦИЯ О СЛАЙДЕ:
{slide_info}

ДОСТУПНЫЕ LAYOUTS:
{layouts_list}

{rules}

Ответь ТОЛЬКО числом - индексом выбранного layout (0-{len(layouts_info)-1})."""
    
    try:
        # Получаем базовый URL если настроен (для прокси или альтернативных API)
        openai_base_url = os.getenv("OPENAI_API_BASE")
        
        if openai_base_url:
            client = AsyncOpenAI(api_key=openai_key, base_url=openai_base_url)
        else:
            client = AsyncOpenAI(api_key=openai_key)
        
        response = await client.chat.completions.create(
            model=os.getenv("MODEL_NAME"),
            messages=[
                {"role": "system", "content": "Ты эксперт по дизайну презентаций. Отвечай только числом."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.3,
            max_tokens=10
        )
        
        # Извлекаем номер layout
        answer = response.choices[0].message.content.strip()
        layout_idx = int(''.join(filter(str.isdigit, answer)))
        
        # Валидация
        if 0 <= layout_idx < len(all_layouts):
            layout_name = all_layouts[layout_idx].name
            return layout_idx
        else:
            return 1 if position > 0 else 0
            
    except Exception as e:
        # Fallback: используем правила на основе названий layouts
        if position == 0:
            # ТИТУЛЬНЫЙ СЛАЙД: ищем layout с "титул" или "title" в названии
            for info in layouts_info:
                if any(word in info["name"].lower() for word in ["title", "титул", "заголов", "обложк"]):
                    return info["index"]
            # Последний fallback для титульника
            return 6 if len(all_layouts) > 6 else 0
        else:
            # КОНТЕНТНЫЙ СЛАЙД: НЕ берем титульные layouts
            # Сначала ищем layout с "контент", "слайд", "логотип" в названии
            for info in layouts_info:
                name_lower = info["name"].lower()
                # Исключаем титульные
                if any(word in name_lower for word in ["title", "титул", "заголов", "обложк"]):
                    continue
                # Берем контентные
                if any(word in name_lower for word in ["content", "контент", "текст", "слайд", "логотип"]):
                    return info["index"]
            # Последний fallback: первый НЕ титульный layout
            for info in layouts_info:
                name_lower = info["name"].lower()
                if not any(word in name_lower for word in ["title", "титул", "заголов", "обложк"]):
                    return info["index"]
            # Совсем последний fallback
            return min(1, len(all_layouts) - 1)


async def generate_pptx_from_slides(
    slides_data: List[Dict[str, Any]],
    template_pptx_bytes: Optional[bytes] = None,
    presentation_title: str = "Presentation"
) -> bytes:
    """
    Генерирует PPTX презентацию напрямую из структурированных данных слайдов
    
    Args:
        slides_data: Список слайдов с полями:
            - title: str
            - position: int
            - content_json: str (JSON with {"blocks": [{"type": "text/list/chart/table", "data": {...}}]})
        template_pptx_bytes: Байты template.pptx для использования layouts
        presentation_title: Название презентации
        
    Логика выбора макетов:
        - ПЕРВЫЙ слайд (position=0): всегда "Титульный_белая тема" (заполняется только заголовок)
        - ОСТАЛЬНЫЕ слайды: выбираются через AI (select_layout_with_ai)
        
    Returns:
        PPTX файл в байтах
    """
    # Загружаем template или создаем новый
    if template_pptx_bytes:
        prs = Presentation(io.BytesIO(template_pptx_bytes))
        template_slide_count = len(prs.slides)
        has_template = True
        
        # Подсчитываем все доступные layouts из всех slide masters
        all_layouts = []
        
        for master_idx, slide_master in enumerate(prs.slide_masters):
            master_layouts = list(slide_master.slide_layouts)
            for layout in master_layouts:
                all_layouts.append(layout)
        
        # Создаем объект для доступа ко всем layouts
        all_layouts_obj = AllLayouts(prs)
        
        # УДАЛЯЕМ ВСЕ слайды из шаблона
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[-1].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[-1]
    else:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        has_template = False
        # Для пустой презентации используем prs.slide_layouts
        all_layouts_obj = prs.slide_layouts
    
    # Функция поиска layout по имени
    def find_layout_by_name(all_layouts, name_pattern: str):
        """Ищет layout по названию (частичное совпадение, без учета регистра)"""
        for idx, layout in enumerate(all_layouts):
            layout_name = layout.name.lower()
            if name_pattern.lower() in layout_name:
                return idx
        return None
    
    # Создаём нужное количество слайдов используя layouts
    num_slides_needed = len(slides_data)
    
    # Сортируем слайды по позиции для правильного порядка
    sorted_slides = sorted(slides_data, key=lambda x: x['position'])
    
    # Определяем layout для каждого слайда
    for i, slide_data in enumerate(sorted_slides):
        position = slide_data['position']
        layout_idx = None
        
        # ПЕРВЫЙ СЛАЙД (титульник) - всегда "Титульный_белая тема"
        # Используем i (индекс) вместо position, так как positions из БД могут начинаться с 1
        if i == 0:
            layout_idx = find_layout_by_name(all_layouts_obj, "Титульный_белая тема")
            if layout_idx is None:
                print(f"⚠️ Layout 'Титульный_белая_тема' не найден, используем первый layout")
                layout_idx = 0
        
        # ВСЕ ОСТАЛЬНЫЕ СЛАЙДЫ (включая последний) - выбираем через AI
        else:
            # Передаем i (индекс в массиве) вместо position из БД для правильной логики титульного/контентного слайда
            layout_idx = await select_layout_with_ai(all_layouts_obj, slide_data, i)
            # Валидация: проверяем что layout_idx в допустимых пределах
            if layout_idx >= len(all_layouts_obj):
                layout_idx = 0
        
        slide_layout = all_layouts_obj[layout_idx]
        prs.slides.add_slide(slide_layout)
        print(f"📄 Слайд {position}: layout '{slide_layout.name}' (idx={layout_idx})")
    
    # Заполняем слайды контентом - используем ТУ ЖЕ sorted_slides что и при создании!
    for idx, slide_data in enumerate(sorted_slides):
        slide = prs.slides[idx]
        position = slide_data.get('position', idx)
        slide_id = slide_data.get('id', 'unknown')
        
        # DEBUG: Логируем какой слайд заполняем
        print(f"\n📝 Заполняем слайд idx={idx}, position={position}, slide_id={slide_id}, title='{slide_data.get('title', '')[:50]}...'")
        print(f"   Layout: {slide.slide_layout.name}")
        
        # Получаем title
        title_text = slide_data.get('title', f'Slide {idx + 1}').strip()
        
        # Парсим content_json
        content_json_str = slide_data.get('content_json', '{}')
        try:
            content_data = json.loads(content_json_str) if isinstance(content_json_str, str) else content_json_str
        except json.JSONDecodeError:
            content_data = {}
        
        blocks = content_data.get('blocks', [])
        
        # DEBUG: Показываем blocks для этого слайда
        print(f"   📦 Blocks: {len(blocks)}")
        for b_idx, b in enumerate(blocks):
            b_type = b.get('type', 'unknown')
            b_data_keys = list(b.get('data', {}).keys())
            print(f"      Block {b_idx+1}: type={b_type}, data_keys={b_data_keys}")
        
        # Заполняем заголовок
        title_placeholder = None
        content_placeholder = None
        
        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type in [1, 3]:  # Title or Center Title
                    title_placeholder = shape
                elif ph_type == 2:  # Body/Content
                    content_placeholder = shape
        
        # Заполняем title placeholder ТОЛЬКО если он есть
        if title_placeholder:
            try:
                title_frame = title_placeholder.text_frame
                title_frame.clear()
                # Truncate заголовок чтобы влез
                truncated_title = _truncate_text_to_fit(title_text, max_length=100)
                title_frame.text = truncated_title.strip()
                
                # Подгоняем размер шрифта заголовка под длину
                title_length = len(truncated_title)
                if title_length > 70:
                    title_font_size = Pt(24)
                elif title_length > 50:
                    title_font_size = Pt(28)
                elif title_length > 30:
                    title_font_size = Pt(32)
                else:
                    title_font_size = Pt(36)
                
                for paragraph in title_frame.paragraphs:
                    paragraph.font.size = title_font_size
                    paragraph.font.bold = True
            except:
                pass  # Если не удалось - пропускаем
        
        # ПЕРВЫЙ СЛАЙД (титульник) - только заголовок, контент не трогаем
        if idx == 0:  # Используем idx вместо position
            continue  # Пропускаем заполнение контента для титульного слайда
        
        # ВСЕ ОСТАЛЬНЫЕ СЛАЙДЫ (включая последний) - заполняем контентом
        # Обрабатываем контентные блоки
        content_was_added = False
        
        if blocks:
            # Очищаем content placeholder только если он есть
            if content_placeholder:
                try:
                    content_frame = content_placeholder.text_frame
                    content_frame.clear()
                except:
                    pass  # Если не удалось очистить - продолжаем
            content_was_added = await _add_blocks_to_slide(slide, blocks, content_placeholder)
        else:
            # Если нет блоков, добавляем description как текст в content placeholder (если есть)
            description = slide_data.get('description', '').strip()
            if description and content_placeholder:
                try:
                    content_frame = content_placeholder.text_frame
                    content_frame.clear()
                    # Truncate description чтобы влез
                    truncated_desc = _truncate_text_to_fit(description, max_length=800)
                    content_frame.text = truncated_desc.strip()
                    
                    # Подгоняем размер шрифта под длину description
                    desc_length = len(truncated_desc)
                    if desc_length > 1500:
                        desc_font_size = Pt(10)
                    elif desc_length > 800:
                        desc_font_size = Pt(11)
                    elif desc_length > 500:
                        desc_font_size = Pt(12)
                    else:
                        desc_font_size = Pt(14)
                    
                    for paragraph in content_frame.paragraphs:
                        paragraph.font.size = desc_font_size
                    content_was_added = True
                except:
                    pass  # Если не удалось - пропускаем
        
        # FALLBACK: если контент НЕ добавлен - пересоздаем слайд с layout "Слайд с заголовком"
        # и добавляем текстовое поле для контента
        if not content_was_added and idx > 0:  # Не для титульного
            print(f"⚠️ Слайд {position} пустой, применяем fallback логику")
            
            # Собираем текст для fallback
            fallback_text = None
            
            # Сначала пробуем взять текст из blocks
            if blocks:
                for block in blocks:
                    if block.get('type') == 'text' and block.get('data', {}).get('text'):
                        fallback_text = block['data']['text'].strip()
                        break
                    elif block.get('type') == 'list' and block.get('data', {}).get('items'):
                        items = block['data']['items']
                        fallback_text = '\n'.join([f"• {item.strip()}" for item in items[:10] if item.strip()])
                        break
            
            # Если нет текста в blocks - используем description
            if not fallback_text and slide_data.get('description'):
                fallback_text = slide_data['description'].strip()
            
            if fallback_text:
                # ПРОБУЕМ ВАРИАНТ 1: Добавить в существующие body placeholders
                added = False
                for shape in slide.shapes:
                    if shape.is_placeholder:
                        ph_type = shape.placeholder_format.type
                        # КРИТИЧНО: НЕ используем title placeholders (1, 3)!
                        if ph_type in [1, 3]:
                            continue
                        try:
                            text_frame = shape.text_frame
                            text_frame.clear()
                            text_frame.text = _truncate_text_to_fit(fallback_text, 1000).strip()
                            for p in text_frame.paragraphs:
                                p.font.size = Pt(12)
                            print(f"  ✅ Добавлен fallback текст в placeholder (type={ph_type})")
                            added = True
                            break
                        except:
                            continue
                
                # ВАРИАНТ 2: Если не получилось - создаем текстовое поле вручную
                if not added:
                    print(f"  📝 Создаем текстовое поле вручную для fallback контента")
                    try:
                        # Создаем текстовое поле под заголовком
                        # Стандартные координаты для контентной области
                        left = Inches(0.5)
                        top = Inches(1.5)  # Под заголовком
                        width = Inches(9.0)  # Почти вся ширина слайда
                        height = Inches(5.0)  # Основная область
                        
                        textbox = slide.shapes.add_textbox(left, top, width, height)
                        text_frame = textbox.text_frame
                        text_frame.word_wrap = True
                        text_frame.text = _truncate_text_to_fit(fallback_text, 1000).strip()
                        
                        # Форматирование
                        for paragraph in text_frame.paragraphs:
                            paragraph.font.size = Pt(14)
                            paragraph.space_after = Pt(12)
                        
                        print(f"  ✅ Текстовое поле создано успешно")
                        added = True
                    except Exception as e:
                        print(f"  ❌ Ошибка создания текстового поля: {e}")
                
                if not added:
                    print(f"  ⚠️ Не удалось добавить даже минимальный контент")
            else:
                print(f"  ⚠️ Нет текста для добавления (ни blocks, ни description)")
    
    # Финальная проверка: убеждаемся что количество слайдов соответствует ожидаемому
    expected_slides = len(slides_data)
    actual_slides = len(prs.slides)
    
    if actual_slides != expected_slides:
        print(f"⚠️ ВНИМАНИЕ: Ожидалось {expected_slides} слайдов, но в презентации {actual_slides}!")
        print(f"  Это может привести к ошибкам при сохранении")
    else:
        print(f"✅ Проверка пройдена: {actual_slides} слайдов в презентации")
    
    # КРИТИЧНО: Двойное сохранение для очистки мусорных XML файлов
    # python-pptx оставляет дубликаты XML от удаленных слайдов в архиве
    # Перезагрузка презентации очищает эти дубликаты
    print("💾 Сохраняем PPTX и пересобираем архив для удаления дубликатов...")
    
    temp_pptx = io.BytesIO()
    prs.save(temp_pptx)
    temp_pptx.seek(0)
    
    # Перезагружаем презентацию - python-pptx читает только валидные слайды
    print("🔄 Перезагружаем презентацию для очистки мусорных XML...")
    prs_clean = Presentation(temp_pptx)
    
    # Сохраняем заново - теперь архив содержит только актуальные XML
    final_pptx = io.BytesIO()
    prs_clean.save(final_pptx)
    final_pptx.seek(0)
    
    pptx_bytes = final_pptx.getvalue()
    print(f"✅ PPTX generated successfully ({len(pptx_bytes)} bytes)")
    
    return pptx_bytes


def _truncate_text_to_fit(text: str, max_length: int = 500) -> str:
    """
    Обрезает текст если он слишком длинный
    
    Args:
        text: Исходный текст
        max_length: Максимальная длина
    
    Returns:
        Обрезанный текст с "..." в конце если был обрезан
    """
    if len(text) <= max_length:
        return text
    return text[:max_length-3] + "..."


def _auto_shrink_text_if_needed(text_frame, max_font_size: int = 14, min_font_size: int = 8) -> int:
    """
    Автоматически уменьшает размер шрифта если текст не влезает в placeholder
    
    Args:
        text_frame: TextFrame object из PowerPoint shape
        max_font_size: Максимальный размер шрифта в pt
        min_font_size: Минимальный размер шрифта в pt
        
    Returns:
        Итоговый размер шрифта в пунктах
    """
    from pptx.util import Pt
    
    # Включаем word wrap один раз перед циклом
    text_frame.word_wrap = True
    
    # Начинаем с максимального размера
    current_size = max_font_size
    # Пытаемся заполнить placeholder с текущим размером шрифта
    # Если не влезает, уменьшаем размер
    for attempt in range(max_font_size - min_font_size + 1):
        try:
            # Устанавливаем размер шрифта для всех параграфов
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(current_size)
            # Пытаемся получить высоту текста после установки размера
            try:
                if len(text_frame.paragraphs) == 0:
                    # Если нет txBody, используем минимальный размер
                    current_size -= 1
                    continue
                # Проверяем по количеству строк в параграфах
                # Это более надежно чем прямое получение высоты
                total_lines = 0
                for paragraph in text_frame.paragraphs:
                    # Приблизительно: один параграф = одна строка + переносы
                    line_count = len(paragraph.text) // 50 + 1  # ~50 символов на строку при размере 10pt
                    total_lines += line_count
                
                # Допустимое количество строк зависит от placeholder (примерно 10-15 строк на стандартный placeholder)
                max_lines = 12  # Conservative estimate
                # Если количество строк в пределах нормы, выходим
                if total_lines <= max_lines:
                    return current_size
                
                # Иначе - уменьшаем размер на 1pt и пытаемся снова
                current_size -= 1
                
            except Exception as e:
                # Если произошла ошибка при проверке, продолжаем уменьшать
                current_size -= 1
                continue
            
        except Exception as e:
            # Если произошла ошибка при установке размера, продолжаем
            current_size -= 1
            continue
    
    # Если ничего не подошло, используем минимальный размер
    print(f"       ⚠️ Текст очень длинный, используем минимальный размер {min_font_size}pt")
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(min_font_size)
    
    return min_font_size


def _classify_placeholders(slide):
    """
    Классифицирует все placeholders на слайде по их назначению
    
    Returns:
        dict: {
            'title': placeholder для заголовка,
            'text': список placeholders для текста,
            'chart': список placeholders для графиков,
            'table': список placeholders для таблиц,
            'picture': список placeholders для картинок,
            'object': список универсальных object placeholders
        }
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
    
    classified = {
        'title': None,
        'text': [],
        'chart': [],
        'table': [],
        'picture': [],
        'object': []
    }
    
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        
        ph_type = shape.placeholder_format.type
        ph_name = shape.name.lower() if hasattr(shape, 'name') else ''
        
        # Title placeholders
        if ph_type in [1, 3]:  # TITLE, CENTER_TITLE
            if classified['title'] is None:
                classified['title'] = shape
        
        # Chart placeholders
        elif ph_type == 14:  # CHART
            classified['chart'].append(shape)
        elif 'диаграмм' in ph_name or 'chart' in ph_name or 'график' in ph_name:
            classified['chart'].append(shape)
        
        # Table placeholders
        elif ph_type == 8:  # TABLE
            classified['table'].append(shape)
        elif 'таблиц' in ph_name or 'table' in ph_name:
            classified['table'].append(shape)
        
        # Picture placeholders
        elif ph_type == 13:  # PICTURE
            classified['picture'].append(shape)
        elif 'изображ' in ph_name or 'picture' in ph_name or 'картинк' in ph_name or 'фото' in ph_name or 'рисун' in ph_name:
            classified['picture'].append(shape)
        
        # Text/Body placeholders (исключаем подзаголовки, footer и заголовки!)
        elif ph_type == 2:  # BODY
            # КРИТИЧНО: Исключаем ЗАГОЛОВКИ по названию (на русском и английском)
            if any(word in ph_name for word in ['заголов', 'title', 'heading', 'header']):
                continue  # Это заголовок, НЕ контентный placeholder!
            
            # Исключаем ПОДЗАГОЛОВКИ по названию
            if 'подзаголов' in ph_name or 'subtitle' in ph_name:
                continue  # Явный подзаголовок - пропускаем
            
            # Получаем размеры слайда
            slide_height = Inches(7.5)
            try:
                if hasattr(slide, 'shapes') and len(slide.shapes) > 0:
                    slide_height = slide.shapes[0].height
            except:
                pass
            
            # Исключаем footer (самый низ слайда)
            if shape.top > slide_height * 0.90:
                continue  # Пропускаем footer
            
            # Исключаем ПОДЗАГОЛОВКИ по размеру и позиции:
            # Это маленькие text placeholders в верхней части слайда
            # Признаки подзаголовка:
            # 1. Находится высоко (в первых 30% слайда после заголовка)
            # 2. Маленький по высоте (< 1 дюйм)
            is_in_subtitle_zone = shape.top < Inches(2.5)  # Верхняя зона
            is_small = shape.height < Inches(1.0)  # Маленький
            
            if is_in_subtitle_zone and is_small:
                # Это скорее всего подзаголовок - пропускаем
                continue
            
            # Берем только нормальные контентные placeholders
            # (большие или расположенные ниже)
            classified['text'].append(shape)
        
        # Object placeholders (универсальные)
        elif ph_type == 7:  # OBJECT
            classified['object'].append(shape)
    
    # Сортируем text placeholders по позиции (сверху вниз)
    # Это гарантирует что сначала заполним подзаголовки, потом основной контент
    classified['text'].sort(key=lambda s: s.top)
    
    return classified


async def _add_blocks_to_slide(slide, blocks: List[Dict[str, Any]], content_placeholder) -> bool:
    """
    Добавляет контентные блоки на слайд, используя подходящие placeholders из template
    
    Поддерживаемые типы блоков:
    - text: простой текст → вставляется в text placeholder
    - list: маркированный список → вставляется в text placeholder  
    - chart: диаграмма (bar/line/pie) → вставляется в chart/object placeholder, название в title placeholder
    - table: таблица → вставляется в table/object placeholder
    - image: изображение → генерируется с помощью Kandinsky и вставляется в picture placeholder
    
    Логика выбора placeholders:
    - Сначала ищем специализированные placeholders (text, chart, table, picture)
    - Для chart/table используем object placeholder как fallback
    - Исключаем подзаголовки (маленькие верхние) и footer (нижние)
    - Если placeholder нет - НЕ вставляем контент (пропускаем блок)
    
    Автоподгонка контента:
    - Динамический размер шрифта в зависимости от объема контента
    - Text: 10-14pt (больше текста = меньше шрифт)
    - List: 10-12pt (больше пунктов = меньше шрифт)
    - Table: 8-10pt (больше ячеек = меньше шрифт)
    - Truncate если контент все равно не влезает
    
    Returns:
        True если хотя бы один блок был успешно добавлен, False иначе
    """
    # Классифицируем placeholders
    placeholders = _classify_placeholders(slide)
    
    # DEBUG: Логируем доступные placeholders
    print(f"     📊 Доступно placeholders: text={len(placeholders['text'])}, chart={len(placeholders['chart'])}, "
          f"table={len(placeholders['table'])}, picture={len(placeholders['picture'])}, object={len(placeholders['object'])}")
    
    # Находим title placeholder для возможности вставки названия диаграммы
    title_placeholder = placeholders.get('title')
    
    # Счётчики использованных placeholders
    used_text_idx = 0
    used_chart_idx = 0
    used_table_idx = 0
    used_picture_idx = 0
    
    # Флаг успешного добавления контента
    content_added = False
    
    for block_idx, block in enumerate(blocks):
        block_type = block.get('type', 'text')
        block_data = block.get('data', {})
        print(f"     🔹 Обрабатываем блок {block_idx+1}/{len(blocks)}: type={block_type}")
        
        if block_type == 'text':
            text_content = block_data.get('text', '').strip()
            if not text_content:
                continue
            
            # Используем text placeholder ТОЛЬКО если он есть
            if used_text_idx < len(placeholders['text']):
                text_ph = placeholders['text'][used_text_idx]
                used_text_idx += 1
                try:
                    text_frame = text_ph.text_frame
                    text_frame.clear()
                    text_frame.word_wrap = True
                    # Truncate текст чтобы влез
                    truncated_text = _truncate_text_to_fit(text_content, max_length=800)
                    text_frame.text = truncated_text.strip()
                    
                    # Подгоняем размер шрифта под длину текста (стартовое значение)
                    text_length = len(truncated_text)
                    if text_length > 500:
                        initial_font_size = 10
                    elif text_length > 300:
                        initial_font_size = 11
                    elif text_length > 150:
                        initial_font_size = 12
                    else:
                        initial_font_size = 14
                    
                    # Автоматически уменьшаем размер если текст не влезает
                    final_font_size = _auto_shrink_text_if_needed(
                        text_frame
                    )
                    
                    content_added = True  # Успешно добавили текст
                    print(f"       ✅ Text добавлен в placeholder {used_text_idx-1} (font: {final_font_size}pt)")
                except Exception as e:
                    print(f"       ❌ Ошибка добавления text: {e}")
            else:
                print(f"       ⚠️ Нет text placeholder (used={used_text_idx}, total={len(placeholders['text'])})")
        
        elif block_type == 'list':
            items = block_data.get('items', [])
            if not items:
                continue
            
            # Strip каждый item в списке
            items = [item.strip() for item in items if isinstance(item, str) and item.strip()]
            
            # Используем text placeholder ТОЛЬКО если он есть
            if used_text_idx < len(placeholders['text']):
                text_ph = placeholders['text'][used_text_idx]
                used_text_idx += 1
                try:
                    text_frame = text_ph.text_frame
                    # КРИТИЧНО: ОЧИЩАЕМ placeholder перед добавлением списка!
                    # Каждый блок (text или list) должен быть в СВОЕМ placeholder
                    text_frame.clear()
                    text_frame.word_wrap = True
                    
                    # Ограничиваем количество items чтобы влезло
                    max_items = 10
                    items_to_show = items[:max_items]
                    
                    # Подгоняем размер шрифта под количество пунктов (стартовое значение)
                    num_items = len(items_to_show)
                    if num_items > 8:
                        initial_font_size = 10
                        max_item_length = 100
                    elif num_items > 5:
                        initial_font_size = 11
                        max_item_length = 120
                    else:
                        initial_font_size = 12
                        max_item_length = 150
                    
                    # Добавляем items как новые параграфы в конец
                    for i, item in enumerate(items_to_show):
                        p = text_frame.add_paragraph()
                        # Truncate каждый item и убираем лишние пробелы
                        truncated_item = _truncate_text_to_fit(item, max_length=max_item_length).strip()
                        p.text = f"• {truncated_item}"
                        p.font.size = Pt(initial_font_size)
                        p.space_after = Pt(4)  # Небольшой отступ между пунктами
                    
                    # Если было больше items - добавляем "..."
                    if len(items) > max_items:
                        p = text_frame.add_paragraph()
                        p.text = "..."
                        p.font.size = Pt(initial_font_size)
                    
                    # Автоматически уменьшаем размер если список не влезает
                    final_font_size = _auto_shrink_text_if_needed(
                        text_frame
                    )
                    
                    content_added = True  # Успешно добавили список
                    print(f"       ✅ List добавлен в placeholder {used_text_idx-1} ({len(items_to_show)} items, font: {final_font_size}pt)")
                except Exception as e:
                    print(f"       ❌ Ошибка добавления list: {e}")
            else:
                print(f"       ⚠️ Нет text placeholder для list (used={used_text_idx}, total={len(placeholders['text'])})")
        
        elif block_type == 'chart':
            chart_type = block_data.get('chart_type', 'bar')
            chart_data_raw = block_data.get('data', {})
            chart_title = block_data.get('title', '').strip()  # Название диаграммы
            
            # Если есть название диаграммы - вставляем его в title placeholder
            if chart_title and title_placeholder:
                try:
                    title_frame = title_placeholder.text_frame
                    title_frame.clear()
                    title_frame.text = chart_title.strip()
                    # Устанавливаем форматирование
                    for paragraph in title_frame.paragraphs:
                        paragraph.font.size = Pt(28)
                        paragraph.font.bold = True
                except:
                    pass
            
            # Используем chart placeholder или object placeholder ТОЛЬКО если он есть
            chart_ph = None
            if used_chart_idx < len(placeholders['chart']):
                chart_ph = placeholders['chart'][used_chart_idx]
                used_chart_idx += 1
            elif used_chart_idx < len(placeholders['chart']) + len(placeholders['object']):
                obj_idx = used_chart_idx - len(placeholders['chart'])
                if obj_idx < len(placeholders['object']):
                    chart_ph = placeholders['object'][obj_idx]
                    used_chart_idx += 1
            
            # Вставляем ТОЛЬКО если есть placeholder
            if chart_ph:
                _add_chart_to_slide(
                    slide,
                    chart_type,
                    chart_data_raw,
                    chart_ph.left,
                    chart_ph.top,
                    chart_ph.width,
                    chart_ph.height
                )
                content_added = True  # Успешно добавили график
            # Если placeholder нет - НЕ вставляем (пропускаем блок)
        
        elif block_type == 'table':
            table_data = block_data.get('data', [])
            if not table_data or len(table_data) == 0:
                continue
            
            # Используем table placeholder или object placeholder ТОЛЬКО если он есть
            table_ph = None
            if used_table_idx < len(placeholders['table']):
                table_ph = placeholders['table'][used_table_idx]
                used_table_idx += 1
            elif used_table_idx < len(placeholders['table']) + len(placeholders['object']):
                obj_idx = used_table_idx - len(placeholders['table'])
                if obj_idx < len(placeholders['object']):
                    table_ph = placeholders['object'][obj_idx]
                    used_table_idx += 1
            
            # Вставляем ТОЛЬКО если есть placeholder
            if table_ph:
                # Ограничиваем размер таблицы чтобы влезла
                max_rows = 8  # Максимум 8 строк (включая заголовок)
                max_cols = 6  # Максимум 6 столбцов
                
                # Обрезаем таблицу до нужного размера
                table_data_truncated = table_data[:max_rows]
                for i in range(len(table_data_truncated)):
                    if len(table_data_truncated[i]) > max_cols:
                        table_data_truncated[i] = table_data_truncated[i][:max_cols]
                
                rows = len(table_data_truncated)
                cols = len(table_data_truncated[0]) if rows > 0 else 0
                
                if rows == 0 or cols == 0:
                    continue
                
                try:
                    table_shape = slide.shapes.add_table(
                        rows, cols, 
                        table_ph.left, 
                        table_ph.top, 
                        table_ph.width, 
                        table_ph.height
                    )
                    table = table_shape.table
                    
                    # Подгоняем размер шрифта под размер таблицы (стартовое значение)
                    # Чем больше данных, тем меньше шрифт
                    if rows > 6 or cols > 5:
                        initial_cell_font_size = 8
                        max_cell_length = 30
                    elif rows > 4 or cols > 3:
                        initial_cell_font_size = 9
                        max_cell_length = 40
                    else:
                        initial_cell_font_size = 10
                        max_cell_length = 50
                    
                    # Заполняем таблицу данными (truncate текст в ячейках)
                    for i, row_data in enumerate(table_data_truncated):
                        for j, cell_value in enumerate(row_data):
                            if j < len(table.columns) and i < len(table.rows):
                                cell = table.cell(i, j)
                                # Truncate текст в ячейке
                                cell_text = str(cell_value)
                                truncated_cell = _truncate_text_to_fit(cell_text, max_length=max_cell_length)
                                cell.text = truncated_cell
                                
                                # Устанавливаем размер шрифта (стартовый)
                                for paragraph in cell.text_frame.paragraphs:
                                    paragraph.font.size = Pt(initial_cell_font_size)
                                    # Заголовок жирным
                                    if i == 0:
                                        paragraph.font.bold = True
                    
                    # Автоматически уменьшаем размер шрифта если таблица не влезает
                    # Проверяем каждую ячейку и уменьшаем шрифт если нужно
                    for row in table.rows:
                        for cell in row.cells:
                            final_size = _auto_shrink_text_if_needed(
                                cell.text_frame
                            )
                    
                    content_added = True  # Успешно добавили таблицу
                except:
                    pass  # Если не удалось создать таблицу - пропускаем
            # Если placeholder нет - НЕ вставляем (пропускаем блок)
        
        elif block_type == 'image':
            # Для изображений генерируем изображение с помощью Kandinsky
            image_description = block_data.get('description', '').strip()
            
            if not image_description:
                print("⚠️ Пропущен image блок: нет описания")
                continue
            
            # Проверяем доступность Kandinsky
            if not KANDINSKY_AVAILABLE:
                print("⚠️ Kandinsky недоступен, пропускаем image блок")
                continue
            
            # Находим picture placeholder
            if used_picture_idx < len(placeholders['picture']):
                pic_ph = placeholders['picture'][used_picture_idx]
                used_picture_idx += 1
                
                try:
                    print(f"🎨 Генерация изображения: '{image_description[:100]}...'")
                    
                    # Получаем размеры placeholder в EMU (English Metric Units)
                    left = pic_ph.left
                    top = pic_ph.top
                    width_emu = pic_ph.width
                    height_emu = pic_ph.height
                    
                    print(f"   📐 Размеры placeholder (EMU): width={width_emu}, height={height_emu}")
                    
                    # Конвертируем EMU в пиксели для Kandinsky API
                    # 1 дюйм = 914400 EMU, 1 дюйм = 96 пиксели (стандартный DPI)
                    EMU_PER_INCH = 914400
                    PIXELS_PER_INCH = 96
                    
                    # Вычисляем размеры в пиксельях
                    width_px = int((width_emu / EMU_PER_INCH) * PIXELS_PER_INCH)
                    height_px = int((height_emu / EMU_PER_INCH) * PIXELS_PER_INCH)
                    
                    print(f"   📐 Размеры после конверсии: width={width_px}px, height={height_px}px")
                    
                    # Убеждаемся что размеры кратны 8 (требование Kandinsky API)
                    # и не меньше минимума
                    MIN_SIZE = 256  # Минимальный размер для Kandinsky
                    
                    # Округляем до ближайшего кратного 8
                    if width_px > 0:
                        width_px = max(MIN_SIZE, (width_px // 8) * 8)
                    else:
                        width_px = 512
                    
                    if height_px > 0:
                        height_px = max(MIN_SIZE, (height_px // 8) * 8)
                    else:
                        height_px = 512
                    
                    print(f"   📐 Размеры после округления: {width_px} x {height_px} px (выровнены к 8)")
                    
                    # Генерируем изображение с помощью Kandinsky
                    # Используем размеры из placeholder
                    print(f"   🔄 Вызываем generate_image_with_retry(width={width_px}, height={height_px})")
                    image_bytes = await generate_image_with_retry(
                        description=image_description,
                        style="DEFAULT",
                        art_gpt=True,
                        width=width_px,
                        height=height_px,
                        max_retries=2
                    )
                    
                    if image_bytes:
                        # ВАЖНО: Удаляем placeholder перед вставкой изображения
                        sp = pic_ph.element
                        sp.getparent().remove(sp)
                        
                        # Вставляем изображение на место placeholder с размерами из placeholder
                        image_stream = io.BytesIO(image_bytes)
                        picture = slide.shapes.add_picture(
                            image_stream,
                            left,
                            top,
                            width=width_emu,
                            height=height_emu
                        )
                        
                        print(f"✅ Изображение добавлено на слайд (сгенерировано {width_px}x{height_px}px, размещено {width_emu}x{height_emu}EMU)")
                        content_added = True
                    else:
                        print(f"⚠️ Не удалось сгенерировать изображение (generate_image_with_retry вернул None)")
                        
                except Exception as e:
                    import traceback
                    print(f"⚠️ Ошибка при добавлении изображения: {str(e)}")
                    print(f"   Traceback: {traceback.format_exc()}")
            # Если placeholder нет - НЕ вставляем (пропускаем блок)
    
    print(f"     📌 _add_blocks_to_slide возвращает: {content_added}")
    return content_added


def _add_chart_to_slide(
    slide,
    chart_type: str,
    chart_data: Dict[str, Any],
    left,
    top,
    width,
    height
):
    """
    Добавляет диаграмму на слайд
    
    chart_data format:
    {
        "categories": ["Category 1", "Category 2", ...],
        "series": [
            {"name": "Series 1", "values": [10, 20, 30, ...]},
            {"name": "Series 2", "values": [15, 25, 35, ...]}
        ]
    }
    """
    # Подготавливаем данные
    categories = chart_data.get('categories', [])
    series_list = chart_data.get('series', [])
    
    # КРИТИЧНО: Автоматически конвертируем bar -> line для одной категории
    # Один столбец выглядит плохо, лучше использовать линейный график
    if chart_type.lower() in ['bar', 'column'] and len(categories) == 1:
        print(f"   🔄 Автоконверсия: bar -> line (одна категория '{categories[0]}')")
        chart_type = 'line'
    
    # Определяем тип диаграммы
    chart_type_map = {
        'bar': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'line': XL_CHART_TYPE.LINE,
        'pie': XL_CHART_TYPE.PIE,
    }
    
    xl_chart_type = chart_type_map.get(chart_type.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)
    
    if not categories or not series_list:
        return
    
    chart_data_obj = CategoryChartData()
    chart_data_obj.categories = categories
    
    for series in series_list:
        series_name = series.get('name', 'Series')
        series_values = series.get('values', [])
        chart_data_obj.add_series(series_name, series_values)
    
    # Добавляем диаграмму
    chart = slide.shapes.add_chart(
        xl_chart_type, left, top, width, height, chart_data_obj
    ).chart
    
    # Настройка диаграммы
    chart.has_legend = True
    chart.legend.position = 2  # Right
    chart.legend.include_in_layout = False
