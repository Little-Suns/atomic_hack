"""Извлекает информацию о дизайне из PPTX, включая изображения, и загружает их в S3."""
import io
import uuid
import asyncio
from typing import Dict, List, Any, Optional
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL
from pptx.slide import Slide, SlideMaster, SlideLayout

# Предполагается, что s3_service находится в том же каталоге
from services import s3_service

def rgb_to_hex(rgb_color) -> str:
    """Конвертирует RGB в HEX."""
    try:
        return f"#{rgb_color.r:02x}{rgb_color.g:02x}{rgb_color.b:02x}"
    except Exception:
        return "#000000"

def get_fill_info(fill, part) -> Dict[str, Any]:
    """Извлекает информацию о заливке, включая изображения."""
    if fill.type == MSO_FILL.SOLID:
        return {
            "type": "solid",
            "color": rgb_to_hex(fill.fore_color.rgb)
        }
    if fill.type == MSO_FILL.PICTURE:
        try:
            image_part = part.related_part(fill.image._rId)
            image_bytes = image_part.blob
            content_type = image_part.content_type
            return {
                "type": "image",
                "data": image_bytes,
                "content_type": content_type
            }
        except Exception:
            return {"type": "image_error"}
    # Другие типы заливки (градиент, узор и т.д.) можно добавить здесь
    return {"type": "other"}


async def extract_design_info(pptx_bytes: bytes) -> Dict[str, Any]:
    """
    Извлекает полную информацию о дизайне презентации, загружает изображения в S3
    и возвращает JSON с URL-ами изображений.
    """
    presentation = Presentation(io.BytesIO(pptx_bytes))
    design_info = {
        "metadata": {
            "slide_width": presentation.slide_width.pt,
            "slide_height": presentation.slide_height.pt,
        },
        "layouts": {}
    }

    upload_tasks = []
    image_map = {} # Карта для отслеживания уже загруженных изображений (hash -> s3_key)

    # Обработка мастер-слайдов и их макетов
    for master in presentation.slide_masters:
        for layout in master.slide_layouts:
            layout_name = layout.name
            layout_info = {
                "background": {},
                "shapes": []
            }

            # Фон макета
            bg_fill = get_fill_info(layout.background.fill, layout.part)
            if bg_fill.get("type") == "image":
                image_hash = hash(bg_fill["data"])
                if image_hash not in image_map:
                    ext = bg_fill["content_type"].split('/')[-1]
                    s3_key = f"template_assets/{uuid.uuid4()}.{ext}"
                    image_map[image_hash] = s3_key
                    upload_tasks.append(
                        s3_service.upload_bytes_to_s3(
                            bg_fill["data"], s3_key, bg_fill["content_type"]
                        )
                    )
                layout_info["background"] = {"type": "image", "s3_key": image_map[image_hash]}
            else:
                layout_info["background"] = bg_fill


            # Фигуры на макете
            for shape in layout.shapes:
                if not shape.is_placeholder:
                    shape_info = await process_shape(shape, layout.part, upload_tasks, image_map)
                    if shape_info:
                        layout_info["shapes"].append(shape_info)
            
            design_info["layouts"][layout_name] = layout_info

    # Загружаем все собранные изображения в S3 параллельно
    if upload_tasks:
        await asyncio.gather(*upload_tasks)

    # Генерируем presigned URLs для всех загруженных изображений
    url_map = {}
    for s3_key in image_map.values():
        url_map[s3_key] = await s3_service.generate_presigned_url(s3_key)

    # Подставляем URL в итоговую структуру
    for layout_info in design_info["layouts"].values():
        if layout_info["background"].get("type") == "image":
            s3_key = layout_info["background"]["s3_key"]
            layout_info["background"]["url"] = url_map.get(s3_key)
        
        for shape_info in layout_info["shapes"]:
            if shape_info.get("fill", {}).get("type") == "image":
                s3_key = shape_info["fill"]["s3_key"]
                shape_info["fill"]["url"] = url_map.get(s3_key)

    return design_info


async def process_shape(shape, part, upload_tasks: list, image_map: dict) -> Optional[Dict[str, Any]]:
    """Обрабатывает отдельную фигуру, извлекая ее свойства и изображение."""
    shape_info = {
        "type": "shape",
        "left": shape.left.pt,
        "top": shape.top.pt,
        "width": shape.width.pt,
        "height": shape.height.pt,
        "fill": {}
    }

    # Заливка фигуры
    fill_info = get_fill_info(shape.fill, part)
    if fill_info.get("type") == "image":
        image_hash = hash(fill_info["data"])
        if image_hash not in image_map:
            ext = fill_info["content_type"].split('/')[-1]
            s3_key = f"template_assets/{uuid.uuid4()}.{ext}"
            image_map[image_hash] = s3_key
            upload_tasks.append(
                s3_service.upload_bytes_to_s3(
                    fill_info["data"], s3_key, fill_info["content_type"]
                )
            )
        shape_info["fill"] = {"type": "image", "s3_key": image_map[image_hash]}
    else:
        shape_info["fill"] = fill_info

    # Другие свойства, как текст, можно добавить здесь
    
    return shape_info