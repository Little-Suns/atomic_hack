"""Конвертирует HTML слайды в PPTX презентацию"""
import io
from typing import List, Dict, Optional, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from bs4 import BeautifulSoup
import re


def hex_to_rgb(hex_color: str) -> tuple:
    """Конвертирует HEX цвет в RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def parse_html_content(html: str) -> Dict[str, Any]:
    """
    Парсит HTML контент слайда и извлекает структурированные данные
    
    Args:
        html: HTML строка
        
    Returns:
        Словарь с title, content_blocks
    """
    soup = BeautifulSoup(html, 'html.parser')
    
    # Извлекаем заголовок
    title_elem = soup.find(['h1', 'h2'], class_=lambda x: x and 'title' in x.lower())
    title = title_elem.get_text(strip=True) if title_elem else ""
    
    # Извлекаем контент
    content_blocks = []
    content_div = soup.find('div', class_=lambda x: x and 'content' in x.lower())
    
    if content_div:
        # Обрабатываем параграфы
        for p in content_div.find_all('p'):
            text = p.get_text(strip=True)
            if text:
                content_blocks.append({'type': 'paragraph', 'text': text})
        
        # Обрабатываем списки
        for ul in content_div.find_all(['ul', 'ol']):
            items = [li.get_text(strip=True) for li in ul.find_all('li')]
            if items:
                content_blocks.append({
                    'type': 'list',
                    'ordered': ul.name == 'ol',
                    'items': items
                })
        
        # Обрабатываем заголовки внутри контента
        for h in content_div.find_all(['h2', 'h3', 'h4']):
            text = h.get_text(strip=True)
            if text:
                content_blocks.append({'type': 'heading', 'text': text, 'level': int(h.name[1])})
        
        # Обрабатываем изображения
        for img in content_div.find_all('img'):
            src = img.get('src', '')
            alt = img.get('alt', 'Image')
            content_blocks.append({
                'type': 'image',
                'src': src,
                'alt': alt
            })
    
    return {
        'title': title,
        'content_blocks': content_blocks
    }


def apply_design_colors(prs: Presentation, design: Optional[Dict]) -> tuple:
    """
    Извлекает цвета из дизайна
    
    Returns:
        (title_color, text_color, bg_color) как RGBColor
    """
    if not design or 'colors' not in design:
        return (RGBColor(0, 0, 0), RGBColor(0, 0, 0), RGBColor(255, 255, 255))
    
    colors = design['colors']
    
    # Берем первые доступные цвета
    title_color = RGBColor(*hex_to_rgb(colors[0])) if len(colors) > 0 else RGBColor(0, 0, 0)
    text_color = RGBColor(*hex_to_rgb(colors[1])) if len(colors) > 1 else RGBColor(0, 0, 0)
    bg_color = RGBColor(*hex_to_rgb(colors[2])) if len(colors) > 2 else RGBColor(255, 255, 255)
    
    return (title_color, text_color, bg_color)


async def convert_html_slides_to_pptx(
    slides_data: List[Dict[str, Any]],
    design_template: Optional[Dict] = None,
    presentation_title: str = "Presentation",
    template_pptx_bytes: Optional[bytes] = None
) -> bytes:
    """
    Конвертирует HTML слайды в PPTX презентацию
    
    Args:
        slides_data: Список слайдов с полями: title, html_content, position
        design_template: JSON дизайн шаблона
        presentation_title: Название презентации
        template_pptx_bytes: Байты template.pptx для использования layouts с дизайном
        
    Returns:
        PPTX файл в байтах
    """
    # Если есть template - используем его напрямую
    if template_pptx_bytes:
        prs = Presentation(io.BytesIO(template_pptx_bytes))
        # Сохраняем слайды шаблона для циклического использования
        template_slide_count = len(prs.slides)
        has_template = True
    else:
        # Fallback: создаем новую презентацию
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        template_slide_count = 0
        has_template = False
    
    # Получаем цвета из дизайна
    title_color, text_color, bg_color = apply_design_colors(prs, design_template)
    
    # Если у нас есть шаблон с слайдами, используем их циклически
    # Если слайдов больше чем в шаблоне - дублируем слайды из шаблона
    if has_template and template_slide_count > 0:
        # Удаляем лишние слайды или дублируем недостающие
        num_slides_needed = len(slides_data)
        
        # Если нужно больше слайдов, дублируем из шаблона
        while len(prs.slides) < num_slides_needed:
            # Дублируем слайд из шаблона (циклически)
            source_idx = len(prs.slides) % template_slide_count
            source_layout = prs.slides[source_idx].slide_layout
            prs.slides.add_slide(source_layout)
        
        # Если слайдов больше чем нужно, удаляем лишние
        while len(prs.slides) > num_slides_needed:
            rId = prs.slides._sldIdLst[-1].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[-1]
    
    # Заполняем слайды контентом
    num_slides = len(slides_data)
    for idx, slide_data in enumerate(sorted(slides_data, key=lambda x: x['position'])):
        # Парсим HTML
        parsed = parse_html_content(slide_data['html_content'])
        
        # Если есть шаблон с слайдами, используем существующий слайд
        if has_template and template_slide_count > 0:
            # Используем уже существующий слайд (мы подготовили нужное количество выше)
            slide = prs.slides[idx]
        else:
            # Создаем новый слайд если нет шаблона
            # Выбираем layout
            if idx == 0:
                slide_layout = prs.slide_layouts[0]
            elif idx == num_slides - 1 and len(prs.slide_layouts) > 2:
                slide_layout = prs.slide_layouts[min(2, len(prs.slide_layouts) - 1)]
            else:
                slide_layout = prs.slide_layouts[min(1, len(prs.slide_layouts) - 1)]
            
            slide = prs.slides.add_slide(slide_layout)
        
        # Ищем placeholders для заголовка и контента в layout
        title_placeholder = None
        content_placeholder = None
        
        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                # 1 = TITLE, 2 = BODY/CONTENT, 3 = CENTER_TITLE
                if ph_type in [1, 3]:  # Title or Center Title
                    title_placeholder = shape
                elif ph_type == 2:  # Body/Content
                    content_placeholder = shape
        
        # Заполняем заголовок
        title_text = parsed['title'] or slide_data['title']
        if title_placeholder:
            # Используем placeholder из layout (сохраняет дизайн!)
            title_frame = title_placeholder.text_frame
            title_frame.clear()
            title_frame.text = title_text
        else:
            # Fallback: создаем textbox
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
            )
            title_frame = title_box.text_frame
            title_frame.text = title_text
            title_paragraph = title_frame.paragraphs[0]
            title_paragraph.font.size = Pt(32)
            title_paragraph.font.bold = True
            title_paragraph.font.color.rgb = title_color
        
        # Заполняем контент
        if content_placeholder:
            # Используем placeholder из layout
            content_frame = content_placeholder.text_frame
            content_frame.clear()
        else:
            # Fallback: создаем textbox
            content_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(1.3), Inches(8.6), Inches(3.5)
            )
            content_frame = content_box.text_frame
        
        content_frame.word_wrap = True
        
        # Добавляем контент блоки
        for i, block in enumerate(parsed['content_blocks']):
            if i > 0:
                # Добавляем параграф для следующего блока
                p = content_frame.add_paragraph()
            else:
                p = content_frame.paragraphs[0]
            
            if block['type'] == 'paragraph':
                p.text = block['text']
                p.font.size = Pt(14)
                p.font.color.rgb = text_color
                p.space_after = Pt(12)
            
            elif block['type'] == 'heading':
                p.text = block['text']
                p.font.size = Pt(18 if block['level'] == 2 else 16)
                p.font.bold = True
                p.font.color.rgb = title_color
                p.space_after = Pt(10)
            
            elif block['type'] == 'list':
                for item_text in block['items']:
                    if p.text:  # Если параграф не пустой, создаем новый
                        p = content_frame.add_paragraph()
                    p.text = f"• {item_text}"
                    p.font.size = Pt(14)
                    p.font.color.rgb = text_color
                    p.level = 0
                    p.space_after = Pt(6)
            
            elif block['type'] == 'image':
                # Добавляем placeholder для изображения
                if p.text:  # Если параграф не пустой, создаем новый
                    p = content_frame.add_paragraph()
                p.text = f"[Изображение: {block['alt']}]"
                p.font.size = Pt(12)
                p.font.italic = True
                p.font.color.rgb = RGBColor(128, 128, 128)
                p.space_after = Pt(10)
    
    # Сохраняем в BytesIO
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    
    return pptx_io.getvalue()
