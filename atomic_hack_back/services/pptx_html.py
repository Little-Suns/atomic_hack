from __future__ import annotations

import base64
import html
import io
from typing import Iterable, List

from pptx import Presentation  # type: ignore[import]


def _iter_text_shapes(shapes: Iterable[object]) -> Iterable[object]:
    for shape in shapes:
        if getattr(shape, "has_text_frame", False):
            yield shape
        if hasattr(shape, "shapes"):
            yield from _iter_text_shapes(shape.shapes)


def _iter_pictures(shapes: Iterable[object]) -> Iterable[object]:
    for shape in shapes:
        if getattr(shape, "image", None) is not None:
            yield shape
        if hasattr(shape, "shapes"):
            yield from _iter_pictures(shape.shapes)


def _paragraph_to_html(paragraph) -> str:
    text = "".join(run.text for run in paragraph.runs).strip()
    if not text:
        return ""
    tag = "p"
    if paragraph.level == 0 and paragraph.runs and len(text.split()) <= 12:
        tag = "h2"
    return f'<{tag} class="level-{paragraph.level}">{html.escape(text)}</{tag}>'


def _shape_text_to_html(shape: object) -> List[str]:
    blocks: List[str] = []
    for paragraph in shape.text_frame.paragraphs:
        block = _paragraph_to_html(paragraph)
        if block:
            blocks.append(block)
    if not blocks:
        return []
    return ["<div class=\"text-box\">"] + blocks + ["</div>"]


def _picture_to_html(picture: object) -> str:
    image = getattr(picture, "image")
    if image is None:
        return ""
    image_bytes = image.blob
    encoded = base64.b64encode(image_bytes).decode()
    mime_type = image.content_type or "image/png"
    return f'<img class="picture" alt="" src="data:{mime_type};base64,{encoded}" />'


def pptx_to_html(pptx_bytes: bytes) -> str:
    """Convert PPTX binary content into a simple HTML representation."""
    presentation = Presentation(io.BytesIO(pptx_bytes))
    html_parts: List[str] = ["<div class=\"presentation\">"]

    for index, slide in enumerate(presentation.slides, start=1):
        slide_parts: List[str] = [f'<section class="slide" data-slide="{index}">']

        for shape in _iter_text_shapes(slide.shapes):
            slide_parts.extend(_shape_text_to_html(shape))

        for picture in _iter_pictures(slide.shapes):
            img_html = _picture_to_html(picture)
            if img_html:
                slide_parts.append(img_html)

        slide_parts.append("</section>")
        html_parts.extend(slide_parts)

    html_parts.append("</div>")
    return "\n".join(html_parts)
